"""
Unified Document Parser
Extracts clean text from all supported formats.
Each parser returns ParseResult with text + metadata only — no LLM involved.
"""
import io, json, csv, logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional, Dict, Any

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".csv", ".json", ".docx", ".xlsx", ".pptx"}


@dataclass
class ParseResult:
    text: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    page_count: Optional[int] = None
    error: Optional[str] = None

    @property
    def success(self) -> bool:
        return self.error is None


def parse_document(file_path: str, file_bytes: bytes, filename: str) -> ParseResult:
    """Route to correct parser by extension."""
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return ParseResult(text="", error=f"Unsupported extension: {ext}")

    parsers = {
        ".pdf":  _parse_pdf,
        ".txt":  _parse_txt,
        ".csv":  _parse_csv,
        ".json": _parse_json,
        ".docx": _parse_docx,
        ".xlsx": _parse_xlsx,
        ".pptx": _parse_pptx,
    }
    try:
        return parsers[ext](file_bytes, filename)
    except Exception as e:
        logger.exception(f"Parser failed for {filename}")
        return ParseResult(text="", error=str(e))


# ── Individual parsers ────────────────────────────────────────────────────────

def _parse_pdf(data: bytes, filename: str) -> ParseResult:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        pages = len(pdf.pages)
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
    return ParseResult(
        text="\n\n".join(text_parts),
        page_count=pages,
        metadata={"filename": filename, "format": "pdf"},
    )


def _parse_txt(data: bytes, filename: str) -> ParseResult:
    # Try common encodings
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            text = data.decode(enc)
            return ParseResult(text=text, metadata={"filename": filename, "encoding": enc, "format": "txt"})
        except UnicodeDecodeError:
            continue
    return ParseResult(text="", error="Cannot decode text file")


def _parse_csv(data: bytes, filename: str) -> ParseResult:
    text_io = io.StringIO(data.decode("utf-8", errors="replace"))
    reader = csv.DictReader(text_io)
    rows = list(reader)
    # Convert to readable text representation
    lines = []
    if reader.fieldnames:
        lines.append(" | ".join(reader.fieldnames))
        lines.append("-" * 60)
    for row in rows:
        lines.append(" | ".join(str(v) for v in row.values()))
    return ParseResult(
        text="\n".join(lines),
        metadata={"filename": filename, "format": "csv", "row_count": len(rows)},
    )


def _parse_json(data: bytes, filename: str) -> ParseResult:
    obj = json.loads(data.decode("utf-8"))
    # Flatten to readable text — policies are often nested JSON
    text = json.dumps(obj, indent=2, ensure_ascii=False)
    return ParseResult(text=text, metadata={"filename": filename, "format": "json"})


def _parse_docx(data: bytes, filename: str) -> ParseResult:
    from docx import Document
    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Include table content
    for table in doc.tables:
        for row in table.rows:
            paragraphs.append(" | ".join(cell.text for cell in row.cells))
    return ParseResult(
        text="\n\n".join(paragraphs),
        metadata={"filename": filename, "format": "docx"},
    )


def _parse_xlsx(data: bytes, filename: str) -> ParseResult:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        parts.append(f"=== Sheet: {sheet} ===")
        for row in ws.iter_rows(values_only=True):
            line = " | ".join(str(c) if c is not None else "" for c in row)
            if line.strip("| "):
                parts.append(line)
    return ParseResult(
        text="\n".join(parts),
        metadata={"filename": filename, "format": "xlsx", "sheets": wb.sheetnames},
    )


def _parse_pptx(data: bytes, filename: str) -> ParseResult:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_parts.append(t)
        slides_text.append("\n".join(slide_parts))
    return ParseResult(
        text="\n\n".join(slides_text),
        page_count=len(prs.slides),
        metadata={"filename": filename, "format": "pptx"},
    )